import os
import sys
import tempfile
import io
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
from openai import OpenAI
import PyPDF2
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Set console encoding to UTF-8
if sys.platform == "win32":
    os.system("chcp 65001 > nul")

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt'}

def allowed_file(filename):
    """Check if the uploaded file has an allowed extension."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf_bytes(file_bytes):
    """Extract text from PDF bytes (avoids filesystem locks on Windows)."""
    try:
        with io.BytesIO(file_bytes) as bio:
            pdf_reader = PyPDF2.PdfReader(bio)
            parts = []
            for page in pdf_reader.pages:
                try:
                    extracted = page.extract_text() or ""
                except Exception:
                    extracted = ""
                parts.append(extracted)
        return "\n".join(parts)
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_text_from_docx_bytes(file_bytes):
    """Extract text from DOCX bytes (avoids filesystem locks on Windows)."""
    try:
        with io.BytesIO(file_bytes) as bio:
            doc = Document(bio)
            lines = [p.text for p in doc.paragraphs]
        return "\n".join(lines)
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def extract_text_from_pptx_bytes(file_bytes):
    """Extract text from PPTX bytes (avoids filesystem locks on Windows)."""
    try:
        with io.BytesIO(file_bytes) as bio:
            prs = Presentation(bio)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def extract_text_from_txt_bytes(file_bytes):
    """Extract text from TXT bytes."""
    try:
        return file_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        return f"Error reading TXT: {str(e)}"

def process_uploaded_file(file):
    """Process uploaded file and extract text content (reads into memory to avoid file locks)."""
    if not allowed_file(file.filename):
        return None, "File type not allowed. Please upload PDF, DOCX, PPTX, or TXT files."

    file_extension = file.filename.rsplit('.', 1)[1].lower()

    # Read the uploaded file into memory once
    try:
        file_bytes = file.read()
    except Exception as e:
        return None, f"Failed to read uploaded file: {str(e)}"

    if file_extension == 'pdf':
        text = extract_text_from_pdf_bytes(file_bytes)
    elif file_extension == 'docx':
        text = extract_text_from_docx_bytes(file_bytes)
    elif file_extension == 'pptx':
        text = extract_text_from_pptx_bytes(file_bytes)
    elif file_extension == 'txt':
        text = extract_text_from_txt_bytes(file_bytes)
    else:
        text = "Unsupported file type"

    return text, None

@app.route('/generate-study-guide', methods=['POST'])
def generate_study_guide():
    """Generate a study guide using the provided parameters and uploaded files."""
    try:
        # Get form data
        subject = request.form.get('subject')
        topic = request.form.get('topic')
        include_practice_exam = request.form.get('include_practice_exam', 'false').lower() == 'true'
        
        if not subject or not topic:
            return jsonify({'error': 'Subject and topic are required'}), 400
        
        # Process uploaded files
        uploaded_content = ""
        if 'files' in request.files:
            files = request.files.getlist('files')
            for file in files:
                if file and file.filename:
                    text, error = process_uploaded_file(file)
                    if error:
                        return jsonify({'error': error}), 400
                    uploaded_content += f"\n\n--- Content from {file.filename} ---\n{text}"
        
        # If no files uploaded, use a default message
        if not uploaded_content:
            uploaded_content = "No additional materials provided"
        
        # Read the template content
        try:
            with open("video_prompt_template.txt", "r", encoding="utf-8") as file:
                template_content = file.read()
        except FileNotFoundError:
            return jsonify({'error': 'Template file not found'}), 500
        
        # Replace the template variables
        prompt = template_content.replace("{{subject}}", subject)
        prompt = prompt.replace("{{topic}}", topic)
        prompt = prompt.replace("{{notes_or_uploaded_materials}}", uploaded_content)
        prompt = prompt.replace("{{true_or_false}}", str(include_practice_exam).lower())
        
        # Generate study guide using OpenAI
        response = client.responses.create(
            model="gpt-5-mini",
            input=prompt
        )
        
        # Return the generated study guide
        return jsonify({
            'success': True,
            'study_guide': response.output_text,
            'subject': subject,
            'topic': topic,
            'include_practice_exam': include_practice_exam,
            'files_processed': len(request.files.getlist('files')) if 'files' in request.files else 0
        })
        
    except Exception as e:
        return jsonify({'error': f'An error occurred: {str(e)}'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint."""
    return jsonify({'status': 'healthy', 'message': 'Study Buddy API is running'})

@app.route('/', methods=['GET'])
def home():
    """Serve the main web interface."""
    return render_template('index.html')

@app.route('/api', methods=['GET'])
def api_docs():
    """API documentation endpoint."""
    return jsonify({
        'message': 'Study Buddy API',
        'endpoints': {
            'POST /generate-study-guide': 'Generate a study guide with uploaded files',
            'GET /health': 'Health check',
            'GET /': 'Web interface',
            'GET /api': 'This documentation'
        },
        'parameters': {
            'subject': 'Class subject (required)',
            'topic': 'Current topic (required)',
            'include_practice_exam': 'Include practice exam (true/false, default: false)',
            'files': 'Upload files (PDF, DOCX, PPTX, TXT) - optional'
        }
    })

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)